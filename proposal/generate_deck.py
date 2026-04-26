"""
North Star — 10-Slide ROI Deck Generator (python-pptx)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)))
DIAG_DIR = os.path.join(os.path.dirname(OUTPUT), "diagrams")

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
GOLD = RGBColor(0xFF, 0xD2, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xEE, 0x6C, 0x4D)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
DARK = RGBColor(0x2C, 0x3E, 0x50)
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)


def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, top=Inches(0.3), left=Inches(0.5), width=Inches(9), fontsize=32, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fontsize)
    p.font.bold = True
    p.font.color.rgb = color


def add_body_text(slide, text, top=Inches(1.5), left=Inches(0.5), width=Inches(9), fontsize=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fontsize)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_table(slide, data, top=Inches(2), left=Inches(0.5), col_widths=None):
    rows, cols = len(data), len(data[0])
    w = col_widths or [Inches(9/cols)] * cols
    table_shape = slide.shapes.add_table(rows, cols, left, top, sum(w), Inches(0.35 * rows))
    table = table_shape.table
    
    for c, width in enumerate(w):
        table.columns[c].width = int(width)
    
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                p.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT
            p.alignment = PP_ALIGN.LEFT
    return table


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(s1, NAVY)
    add_title_box(s1, "NORTH STAR", top=Inches(2), fontsize=44, color=GOLD)
    add_title_box(s1, "Workload Optimizer", top=Inches(2.8), fontsize=36)
    add_body_text(s1, "Digital Transformation for Enterprise Expense Management\n\nTechnology Strategy & Transformation\nProcess Reengineering | Solution Architecture | Automation",
                  top=Inches(4), fontsize=16)
    add_body_text(s1, "Confidential | Prepared for Leadership Review", top=Inches(6.5), fontsize=11, color=GOLD)

    # Slide 2: Executive Summary
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s2, WHITE)
    add_title_box(s2, "Executive Summary", color=NAVY)
    add_body_text(s2, 
        "PROBLEM: Manual expense reporting costs $620K/year with 9.7-day cycle times\n\n"
        "SOLUTION: AI-powered automation reducing cycle time by 78% (9.7 to 2.1 days)\n\n"
        "IMPACT: $434K annual savings | 70% cost reduction | 3.1x productivity gain\n\n"
        "APPROACH: Phased implementation over 12 weeks using Power Platform + Python\n\n"
        "ROI: 865% return in Year 1 with 1.2-month payback period",
        top=Inches(1.5), fontsize=16, color=DARK)

    # Slide 3: Problem Statement
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s3, WHITE)
    add_title_box(s3, "Problem Statement", color=NAVY)
    add_body_text(s3,
        "The current expense reporting process is manual, error-prone, and slow:\n\n"
        "  - 12-step manual workflow across 4 organizational layers\n"
        "  - Average 9.7-day end-to-end processing time\n"
        "  - 18% data entry error rate requiring rework\n"
        "  - $620,647 annual operational cost\n"
        "  - 81.5% receipt attachment rate (target: 98%)\n"
        "  - Employee satisfaction: LOW (delayed reimbursements)\n\n"
        "Data Source: Analysis of 5,000 expense records over 18 months",
        top=Inches(1.5), fontsize=14, color=DARK)

    # Slide 4: AS-IS Flow
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s4, WHITE)
    add_title_box(s4, "AS-IS Process Flow", color=NAVY)
    as_is_path = os.path.join(DIAG_DIR, "as_is_process_flow.png")
    if os.path.exists(as_is_path):
        s4.shapes.add_picture(as_is_path, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.5))

    # Slide 5: Bottleneck Analysis
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s5, WHITE)
    add_title_box(s5, "Bottleneck Analysis", color=NAVY)
    add_table(s5, [
        ['Bottleneck', 'Root Cause', 'Time Lost', 'Annual Cost'],
        ['Manual Data Entry', 'No OCR/auto-import', '33 min/report', '$124,740'],
        ['Approval Queue Delays', 'Email routing, no SLA', '72+ hrs avg wait', '$233,100'],
        ['Finance Reconciliation', 'Manual validation, batch', '40 min + 5-day batch', '$184,800'],
    ], top=Inches(1.5), col_widths=[Inches(2.5), Inches(2.5), Inches(2), Inches(2)])

    # Slide 6: TO-BE Flow
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s6, WHITE)
    add_title_box(s6, "TO-BE Process Flow (Automated)", color=NAVY)
    to_be_path = os.path.join(DIAG_DIR, "to_be_process_flow.png")
    if os.path.exists(to_be_path):
        s6.shapes.add_picture(to_be_path, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.5))

    # Slide 7: Solution Architecture
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s7, WHITE)
    add_title_box(s7, "Solution Architecture", color=NAVY)
    arch_path = os.path.join(DIAG_DIR, "architecture_diagram.png")
    if os.path.exists(arch_path):
        s7.shapes.add_picture(arch_path, Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.5))

    # Slide 8: ROI Analysis
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s8, WHITE)
    add_title_box(s8, "ROI Analysis", color=NAVY)
    add_table(s8, [
        ['Metric', 'Current', 'Future', 'Improvement'],
        ['Annual Operating Cost', '$620,647', '$186,194', '70% reduction'],
        ['Cycle Time', '9.7 days', '2.1 days', '78% faster'],
        ['Reports per FTE/day', '8', '25', '3.1x increase'],
        ['Error Rate', '18%', '5%', '72% reduction'],
        ['Annual Savings', '—', '$434,453', '—'],
        ['ROI (Year 1)', '—', '865%', '—'],
        ['Payback Period', '—', '1.2 months', '—'],
    ], top=Inches(1.3), col_widths=[Inches(2.5), Inches(2), Inches(2), Inches(2.5)])

    # Slide 9: Risk Register
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s9, WHITE)
    add_title_box(s9, "Risk Register", color=NAVY)
    add_table(s9, [
        ['Risk', 'Likelihood', 'Impact', 'Mitigation', 'Status'],
        ['User adoption resistance', 'Medium', 'High', 'Change mgmt + training', 'AMBER'],
        ['Data migration errors', 'Low', 'High', 'Phased rollout + validation', 'GREEN'],
        ['Integration complexity', 'Medium', 'Medium', 'API-first architecture', 'AMBER'],
        ['Budget overrun', 'Low', 'Medium', 'Agile sprints + checkpoints', 'GREEN'],
        ['Vendor lock-in', 'Low', 'Low', 'Modular design + standards', 'GREEN'],
    ], top=Inches(1.3), col_widths=[Inches(2.2), Inches(1.2), Inches(1), Inches(2.8), Inches(1.2)])

    # Slide 10: Implementation Roadmap
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s10, WHITE)
    add_title_box(s10, "Implementation Roadmap", color=NAVY)
    add_table(s10, [
        ['Phase', 'Weeks', 'Deliverables', 'Owner'],
        ['1. Discovery & Design', '1-2', 'AS-IS map, requirements, architecture', 'Consultant'],
        ['2. Data Engineering', '3-4', 'ETL pipeline, database, analytics', 'Data Engineer'],
        ['3. Dashboard & Reporting', '5-6', 'Power BI dashboard, KPI framework', 'BI Developer'],
        ['4. Automation Build', '7-9', 'Power Automate flows, AI detection', 'Automation Dev'],
        ['5. Testing & Training', '10-11', 'UAT, change management, documentation', 'QA + PM'],
        ['6. Go-Live & Handover', '12', 'Production deployment, support handover', 'Full Team'],
    ], top=Inches(1.3), col_widths=[Inches(2.2), Inches(1), Inches(3.8), Inches(2)])

    # Save
    pptx_path = os.path.join(OUTPUT, "ROI_deck.pptx")
    prs.save(pptx_path)
    print(f"PPTX saved: {pptx_path}")

if __name__ == "__main__":
    build_deck()
